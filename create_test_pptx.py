"""
Generate test PPTX files with varying similarity levels for testing.
Run this once, then delete it.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
import os

output_folder = "pptx_files"
os.makedirs(output_folder, exist_ok=True)


def create_pptx(filename, slides_data):
    prs = Presentation()
    for title, body in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    prs.save(os.path.join(output_folder, filename))
    print(f"Created {filename}")


# ============================================================
# FILE 1: Original - Company Q1 Sales Report
# ============================================================
create_pptx("q1_sales_report.pptx", [
    ("Q1 2025 Sales Report", "Prepared by the Sales Department\nConfidential - Internal Use Only"),
    ("Executive Summary",
     "Total revenue for Q1 2025 reached 2.4 million euro, representing a 12% increase over Q1 2024. "
     "The growth was primarily driven by strong performance in the Benelux region and the successful "
     "launch of our new SaaS product line. Customer acquisition cost decreased by 8% while retention "
     "rates improved to 94%."),
    ("Regional Breakdown",
     "Benelux: 1.1M euro (46% of total) - up 18%\n"
     "DACH: 720K euro (30% of total) - up 9%\n"
     "Nordics: 580K euro (24% of total) - up 5%\n"
     "The Benelux region continues to be our strongest market with Belgium and Netherlands "
     "showing exceptional growth in enterprise accounts."),
    ("Key Wins",
     "Signed 3 new enterprise contracts worth 450K euro combined.\n"
     "Expanded partnership with TechCorp for cloud integration services.\n"
     "Launched self-service portal reducing onboarding time by 40%."),
    ("Q2 Outlook",
     "We project Q2 revenue of 2.6 million euro based on current pipeline. "
     "Key priorities include expanding into the French market and launching version 2.0 "
     "of our analytics dashboard. Hiring plan includes 5 new sales representatives."),
])

# ============================================================
# FILE 2: EXACT DUPLICATE (should be ~1.0 similarity)
# ============================================================
create_pptx("q1_sales_report_copy.pptx", [
    ("Q1 2025 Sales Report", "Prepared by the Sales Department\nConfidential - Internal Use Only"),
    ("Executive Summary",
     "Total revenue for Q1 2025 reached 2.4 million euro, representing a 12% increase over Q1 2024. "
     "The growth was primarily driven by strong performance in the Benelux region and the successful "
     "launch of our new SaaS product line. Customer acquisition cost decreased by 8% while retention "
     "rates improved to 94%."),
    ("Regional Breakdown",
     "Benelux: 1.1M euro (46% of total) - up 18%\n"
     "DACH: 720K euro (30% of total) - up 9%\n"
     "Nordics: 580K euro (24% of total) - up 5%\n"
     "The Benelux region continues to be our strongest market with Belgium and Netherlands "
     "showing exceptional growth in enterprise accounts."),
    ("Key Wins",
     "Signed 3 new enterprise contracts worth 450K euro combined.\n"
     "Expanded partnership with TechCorp for cloud integration services.\n"
     "Launched self-service portal reducing onboarding time by 40%."),
    ("Q2 Outlook",
     "We project Q2 revenue of 2.6 million euro based on current pipeline. "
     "Key priorities include expanding into the French market and launching version 2.0 "
     "of our analytics dashboard. Hiring plan includes 5 new sales representatives."),
])

# ============================================================
# FILE 3: SMALL CHANGES / UPDATED VERSION (should be ~0.85-0.95)
# Same structure, minor number updates and small wording tweaks
# ============================================================
create_pptx("q1_sales_report_v2.pptx", [
    ("Q1 2025 Sales Report - REVISED", "Prepared by the Sales Department\nConfidential - Internal Use Only\nRevision 2 - Updated Figures"),
    ("Executive Summary",
     "Total revenue for Q1 2025 reached 2.5 million euro, representing a 14% increase over Q1 2024. "
     "The growth was primarily driven by strong performance in the Benelux region and the successful "
     "launch of our new SaaS product line. Customer acquisition cost decreased by 10% while retention "
     "rates improved to 95%."),
    ("Regional Breakdown",
     "Benelux: 1.15M euro (46% of total) - up 19%\n"
     "DACH: 750K euro (30% of total) - up 11%\n"
     "Nordics: 600K euro (24% of total) - up 7%\n"
     "The Benelux region continues to be our strongest market with Belgium and Netherlands "
     "showing exceptional growth in enterprise accounts."),
    ("Key Wins",
     "Signed 3 new enterprise contracts worth 480K euro combined.\n"
     "Expanded partnership with TechCorp for cloud integration services.\n"
     "Launched self-service portal reducing onboarding time by 42%.\n"
     "Added: Won innovation award at Brussels Tech Summit."),
    ("Q2 Outlook",
     "We project Q2 revenue of 2.7 million euro based on updated pipeline. "
     "Key priorities include expanding into the French market and launching version 2.0 "
     "of our analytics dashboard. Hiring plan includes 6 new sales representatives and 2 solution architects."),
])

# ============================================================
# FILE 4: SAME TOPIC DIFFERENT QUARTER (should be ~0.7-0.8)
# Similar structure and language but different data entirely
# ============================================================
create_pptx("q2_sales_report.pptx", [
    ("Q2 2025 Sales Report", "Prepared by the Sales Department\nConfidential - Internal Use Only"),
    ("Executive Summary",
     "Total revenue for Q2 2025 reached 2.7 million euro, representing a 15% increase over Q2 2024. "
     "Growth was driven by the successful expansion into France and continued momentum in the Benelux "
     "market. The new analytics dashboard contributed 300K euro in upsell revenue. Customer churn "
     "decreased to an all-time low of 4%."),
    ("Regional Breakdown",
     "Benelux: 1.2M euro (44% of total) - up 15%\n"
     "DACH: 800K euro (30% of total) - up 12%\n"
     "Nordics: 500K euro (19% of total) - up 3%\n"
     "France: 200K euro (7% of total) - new market\n"
     "France launch exceeded expectations with 12 new accounts signed in first quarter of operations."),
    ("Key Wins",
     "Closed mega-deal with EuroBank worth 280K euro annual contract.\n"
     "Analytics dashboard 2.0 launched with 89% adoption rate among existing customers.\n"
     "Opened Paris office with 4 initial team members."),
    ("Q3 Outlook",
     "We target Q3 revenue of 2.9 million euro. "
     "Focus areas include deepening France operations, launching a partner reseller program, "
     "and beginning development of the AI-powered document analysis feature."),
])

# ============================================================
# FILE 5: DIFFERENT TOPIC (should be ~0.3-0.5)
# Completely different subject - IT security policy
# ============================================================
create_pptx("it_security_training.pptx", [
    ("IT Security Awareness Training 2025", "Information Security Department\nMandatory for all employees"),
    ("Password Best Practices",
     "Use a minimum of 14 characters for all passwords. "
     "Enable multi-factor authentication on every account. "
     "Never reuse passwords across different services. "
     "Use a company-approved password manager to store credentials securely. "
     "Change passwords immediately if a breach is suspected."),
    ("Phishing Prevention",
     "Always verify the sender email address before clicking any links. "
     "Look for spelling errors and urgency tactics in suspicious emails. "
     "Never download attachments from unknown senders. "
     "Report suspicious emails to security@company.com immediately. "
     "When in doubt, contact the sender through a separate channel to confirm."),
    ("Data Classification",
     "Public: Marketing materials, blog posts, job listings.\n"
     "Internal: Meeting notes, project plans, internal presentations.\n"
     "Confidential: Customer data, financial reports, HR records.\n"
     "Restricted: Encryption keys, authentication tokens, medical data.\n"
     "Always label documents with the correct classification level."),
    ("Incident Response",
     "If you suspect a security incident, disconnect from the network immediately. "
     "Contact the IT Security team via phone at extension 5555. "
     "Do not attempt to investigate or fix the issue yourself. "
     "Preserve all evidence and document what you observed. "
     "A full incident report must be filed within 24 hours."),
])

# ============================================================
# FILE 6: SLIGHTLY RELATED TOPIC (should be ~0.5-0.65)
# Employee onboarding - mentions some similar business concepts
# ============================================================
create_pptx("employee_onboarding.pptx", [
    ("New Employee Onboarding Guide", "Human Resources Department\nWelcome to the team!"),
    ("Company Overview",
     "We are a leading SaaS company operating across the Benelux, DACH, and Nordic regions. "
     "Our flagship product helps enterprises manage document workflows and analytics. "
     "Founded in 2018, we have grown to over 200 employees with offices in Brussels, "
     "Amsterdam, and Munich. Our revenue exceeded 10 million euro in 2024."),
    ("Your First Week",
     "Day 1: IT setup, security training, and team introductions.\n"
     "Day 2: Product overview and customer persona workshop.\n"
     "Day 3: Department-specific onboarding and tool training.\n"
     "Day 4: Shadow a colleague and observe client meetings.\n"
     "Day 5: Set goals with your manager and review onboarding checklist."),
    ("Tools and Access",
     "Email and calendar: Microsoft 365 suite.\n"
     "Communication: Slack for internal messaging.\n"
     "Project management: Jira for task tracking.\n"
     "Documentation: Confluence for company wiki.\n"
     "CRM: Salesforce for customer relationship management.\n"
     "Request access through the IT self-service portal."),
    ("Key Policies",
     "Work from home: Up to 3 days per week after probation period. "
     "Expense reports must be submitted within 30 days using the finance portal. "
     "Annual leave: 25 days per year plus Belgian public holidays. "
     "Performance reviews are conducted quarterly with your direct manager."),
])

print(f"\nDone! Created 6 test PPTX files in '{output_folder}/'")
print("\nExpected similarity levels:")
print("  q1_sales_report.pptx  vs  q1_sales_report_copy.pptx  → ~1.0  (exact duplicate)")
print("  q1_sales_report.pptx  vs  q1_sales_report_v2.pptx    → ~0.85-0.95 (minor updates)")
print("  q1_sales_report.pptx  vs  q2_sales_report.pptx       → ~0.7-0.8  (same topic, different data)")
print("  q1_sales_report.pptx  vs  employee_onboarding.pptx   → ~0.5-0.65 (loosely related)")
print("  q1_sales_report.pptx  vs  it_security_training.pptx  → ~0.3-0.5  (different topic)")
