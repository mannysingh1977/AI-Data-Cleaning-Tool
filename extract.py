import os
import json
from datetime import datetime
from docx import Document
from pypdf import PdfReader
from pptx import Presentation
from openpyxl import load_workbook

def extract_text_from_word(file_path):
    """Extract all text from a Word document"""
    doc = Document(file_path)
    full_text = []
    
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            full_text.append(paragraph.text)
    
    return "\n".join(full_text)

def extract_text_from_pdf(file_path):
    """Extract all text from a PDF with better formatting"""
    reader = PdfReader(file_path)
    text = ""
    
    for page_num, page in enumerate(reader.pages):
        page_text = page.extract_text()
        if page_text:
            # Add page number marker (optional)
            text += f"\n--- Page {page_num + 1} ---\n"
            
            # Split by newlines and rejoin
            lines = page_text.split('\n')
            for line in lines:
                line = line.strip()
                if line:  # Skip empty lines
                    text += line + "\n"
            
            text += "\n"  # Extra space between pages
    
    return text
def extract_text_from_pptx(file_path):
    """Extract all text from a PowerPoint presentation"""
    prs = Presentation(file_path)
    full_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        slide_text.append(" | ".join(row_text))
        if slide_text:
            full_text.append(f"--- Slide {slide_num} ---")
            full_text.extend(slide_text)
            full_text.append("")

    return "\n".join(full_text)

def extract_text_from_xlsx(file_path):
    """Extract all text from an Excel workbook"""
    wb = load_workbook(file_path, data_only=True)
    full_text = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        sheet_text = []
        for row in sheet.iter_rows(values_only=True):
            row_values = [str(cell) for cell in row if cell is not None]
            if row_values:
                sheet_text.append(" | ".join(row_values))
        if sheet_text:
            full_text.append(f"--- Sheet: {sheet_name} ---")
            full_text.extend(sheet_text)
            full_text.append("")

    return "\n".join(full_text)

def extract_metadata_from_word(file_path):
    """Extract metadata from a Word document"""
    doc = Document(file_path)
    props = doc.core_properties
    return {
        "author": props.author or "Unknown",
        "last_modified_by": props.last_modified_by or "Unknown",
        "created": str(props.created) if props.created else "Unknown",
        "modified": str(props.modified) if props.modified else "Unknown",
        "title": props.title or "Unknown",
        "revision": props.revision or 0,
        "page_count": len(doc.paragraphs),
    }

def extract_metadata_from_pdf(file_path):
    """Extract metadata from a PDF"""
    reader = PdfReader(file_path)
    meta = reader.metadata or {}
    return {
        "author": meta.get("/Author", "Unknown") or "Unknown",
        "creator": meta.get("/Creator", "Unknown") or "Unknown",
        "created": str(meta.get("/CreationDate", "Unknown")),
        "modified": str(meta.get("/ModDate", "Unknown")),
        "title": meta.get("/Title", "Unknown") or "Unknown",
        "page_count": len(reader.pages),
    }

def extract_metadata_from_pptx(file_path):
    """Extract metadata from a PowerPoint presentation"""
    prs = Presentation(file_path)
    props = prs.core_properties
    return {
        "author": props.author or "Unknown",
        "last_modified_by": props.last_modified_by or "Unknown",
        "created": str(props.created) if props.created else "Unknown",
        "modified": str(props.modified) if props.modified else "Unknown",
        "title": props.title or "Unknown",
        "revision": props.revision or 0,
        "slide_count": len(prs.slides),
    }

def extract_metadata_from_xlsx(file_path):
    """Extract metadata from an Excel workbook"""
    wb = load_workbook(file_path, data_only=True)
    props = wb.properties
    return {
        "author": props.creator or "Unknown",
        "last_modified_by": props.lastModifiedBy or "Unknown",
        "created": str(props.created) if props.created else "Unknown",
        "modified": str(props.modified) if props.modified else "Unknown",
        "title": props.title or "Unknown",
        "sheet_count": len(wb.sheetnames),
        "sheet_names": wb.sheetnames,
    }

def save_text_to_file(text, output_path):
    """Save extracted text to a .txt file"""
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(text)
    print(f"   ✓ Saved to {output_path}")

if __name__ == "__main__":
    # Define folders
    word_folder = "word_files"
    pdf_folder = "pdf_files"
    pptx_folder = "pptx_files"
    xlsx_folder = "xlsx_files"
    output_folder = "extracted_text"
    metadata_folder = "metadata"

    # Create folders if they don't exist
    for folder in [word_folder, pdf_folder, pptx_folder, xlsx_folder, output_folder, metadata_folder]:
        if not os.path.exists(folder):
            os.makedirs(folder)
            print(f"Created folder: {folder}")
    
    print("\n" + "="*60)
    print("DOCUMENT TEXT EXTRACTION")
    print("="*60)

    all_metadata = {}

    # Process Word documents
    print(f"\n📄 Processing Word documents from '{word_folder}/'...")
    word_count = 0

    if os.path.exists(word_folder):
        for filename in os.listdir(word_folder):
            if filename.endswith(".docx"):
                word_count += 1
                print(f"\n  [{word_count}] {filename}")

                file_path = os.path.join(word_folder, filename)
                text = extract_text_from_word(file_path)
                metadata = extract_metadata_from_word(file_path)

                base_name = os.path.splitext(filename)[0]
                output_path = os.path.join(output_folder, f"{base_name}.txt")
                save_text_to_file(text, output_path)

                metadata["source_file"] = filename
                metadata["source_type"] = "docx"
                metadata["char_count"] = len(text)
                metadata["word_count"] = len(text.split())
                all_metadata[f"{base_name}.txt"] = metadata

                print(f"      Characters: {len(text)} | Words: {len(text.split())}")
                print(f"      Author: {metadata['author']} | Modified: {metadata['modified']}")

    if word_count == 0:
        print(f"   No Word files found. Add .docx files to '{word_folder}/' folder.")

    # Process PDF documents
    print(f"\n📕 Processing PDF documents from '{pdf_folder}/'...")
    pdf_count = 0

    if os.path.exists(pdf_folder):
        for filename in os.listdir(pdf_folder):
            if filename.endswith(".pdf"):
                pdf_count += 1
                print(f"\n  [{pdf_count}] {filename}")

                file_path = os.path.join(pdf_folder, filename)
                text = extract_text_from_pdf(file_path)
                metadata = extract_metadata_from_pdf(file_path)

                base_name = os.path.splitext(filename)[0]
                output_path = os.path.join(output_folder, f"{base_name}.txt")
                save_text_to_file(text, output_path)

                metadata["source_file"] = filename
                metadata["source_type"] = "pdf"
                metadata["char_count"] = len(text)
                metadata["word_count"] = len(text.split())
                all_metadata[f"{base_name}.txt"] = metadata

                print(f"      Characters: {len(text)} | Words: {len(text.split())} | Pages: {metadata['page_count']}")
                print(f"      Author: {metadata['author']} | Modified: {metadata['modified']}")

    if pdf_count == 0:
        print(f"   No PDF files found. Add .pdf files to '{pdf_folder}/' folder.")

    # Process PowerPoint documents
    print(f"\n📊 Processing PowerPoint documents from '{pptx_folder}/'...")
    pptx_count = 0

    if os.path.exists(pptx_folder):
        for filename in os.listdir(pptx_folder):
            if filename.endswith(".pptx"):
                pptx_count += 1
                print(f"\n  [{pptx_count}] {filename}")

                file_path = os.path.join(pptx_folder, filename)
                text = extract_text_from_pptx(file_path)
                metadata = extract_metadata_from_pptx(file_path)

                base_name = os.path.splitext(filename)[0]
                output_path = os.path.join(output_folder, f"{base_name}.txt")
                save_text_to_file(text, output_path)

                metadata["source_file"] = filename
                metadata["source_type"] = "pptx"
                metadata["char_count"] = len(text)
                metadata["word_count"] = len(text.split())
                all_metadata[f"{base_name}.txt"] = metadata

                print(f"      Characters: {len(text)} | Words: {len(text.split())} | Slides: {metadata['slide_count']}")
                print(f"      Author: {metadata['author']} | Modified: {metadata['modified']}")

    if pptx_count == 0:
        print(f"   No PowerPoint files found. Add .pptx files to '{pptx_folder}/' folder.")

    # Process Excel documents
    print(f"\n📗 Processing Excel documents from '{xlsx_folder}/'...")
    xlsx_count = 0

    if os.path.exists(xlsx_folder):
        for filename in os.listdir(xlsx_folder):
            if filename.endswith(".xlsx"):
                xlsx_count += 1
                print(f"\n  [{xlsx_count}] {filename}")

                file_path = os.path.join(xlsx_folder, filename)
                text = extract_text_from_xlsx(file_path)
                metadata = extract_metadata_from_xlsx(file_path)

                base_name = os.path.splitext(filename)[0]
                output_path = os.path.join(output_folder, f"{base_name}.txt")
                save_text_to_file(text, output_path)

                metadata["source_file"] = filename
                metadata["source_type"] = "xlsx"
                metadata["char_count"] = len(text)
                metadata["word_count"] = len(text.split())
                all_metadata[f"{base_name}.txt"] = metadata

                print(f"      Characters: {len(text)} | Words: {len(text.split())} | Sheets: {metadata['sheet_count']}")
                print(f"      Author: {metadata['author']} | Modified: {metadata['modified']}")

    if xlsx_count == 0:
        print(f"   No Excel files found. Add .xlsx files to '{xlsx_folder}/' folder.")

    # Save metadata
    metadata_path = os.path.join(metadata_folder, "metadata.json")
    with open(metadata_path, "w", encoding="utf-8") as f:
        json.dump(all_metadata, f, indent=2, default=str)
    print(f"\n📋 Metadata saved to {metadata_path}")

    # Summary
    print("\n" + "="*60)
    total = word_count + pdf_count + pptx_count + xlsx_count
    print(f"✅ COMPLETE: Processed {total} documents ({word_count} Word + {pdf_count} PDF + {pptx_count} PowerPoint + {xlsx_count} Excel)")
    print(f"📁 Extracted text saved to '{output_folder}/' folder")
    print("="*60)
